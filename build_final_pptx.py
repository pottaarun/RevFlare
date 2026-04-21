#!/usr/bin/env python3
"""
RevFlare Final Presentation Builder
- Starts a mock API server
- Captures screenshots with Playwright
- Builds a polished PPTX with screenshots + business value
"""
import http.server, json, os, threading, time, signal, sys

os.environ['NODE_TLS_REJECT_UNAUTHORIZED'] = '0'

ROOT = os.path.dirname(os.path.abspath(__file__))
PUBLIC = os.path.join(ROOT, 'public')
SHOTS = os.path.join(ROOT, 'screenshots')
os.makedirs(SHOTS, exist_ok=True)

# ══════════════════════════════════════════════════════════════════
# STEP 1: Mock Server
# ══════════════════════════════════════════════════════════════════
MOCK = {
    '/api/stats': {"totalAccounts":2307,"activeAccounts":20,"avgITSpend":78400,"avgCDNSpend":17200,"totalMRR":131000,"totalAddressable":167000000,"openPipeline":62},
    '/api/platform-stats': {"totalEmails":9,"totalResearch":11,"totalCampaigns":1,"totalUsers":1},
    '/api/me': {"email":"apotta@cloudflare.com","accountCount":2307},
    '/api/gmail/status': {"connected":False},
    '/api/settings/status': {},
    '/api/filters': {"industries":["Telecommunications","Computer Software","Media","Publishing","Financial Services"],"countries":["CA","US","GB"],"segments":["Enterprise"],"statuses":["Active","Paid","Unknown"]},
    '/api/accounts': {"accounts":[
        {"id":1,"account_name":"Apple Canada Inc","industry":"Telecommunications","status":"Unknown","total_it_spend":35300000,"current_monthly_fee":0,"cdn_primary":"Akamai","security_primary":"CrowdStrike","billing_country":"CA","employees":5000},
        {"id":2,"account_name":"Equinix Canada Ltd","industry":"Data Centers","status":"Unknown","total_it_spend":22000000,"current_monthly_fee":0,"cdn_primary":"Amazon CloudFront","security_primary":"Neustar","billing_country":"CA","employees":12000},
        {"id":3,"account_name":"Hotbot Limited","industry":"Publishing","status":"Paid","total_it_spend":16600000,"current_monthly_fee":0,"cdn_primary":"Amazon CloudFront","security_primary":"Palo Alto Networks","billing_country":"CA","employees":800,"last_activity":"2026-01-15"},
        {"id":4,"account_name":"Intuit Canada Ulc","industry":"Computer Software","status":"Unknown","total_it_spend":15100000,"current_monthly_fee":0,"cdn_primary":"Blue Coat Systems","security_primary":"Verisign","billing_country":"CA","employees":15000},
        {"id":5,"account_name":"Shopify Inc","industry":"Computer Software","status":"Paid","total_it_spend":45200000,"current_monthly_fee":5400,"cdn_primary":"Fastly","security_primary":"CrowdStrike","billing_country":"CA","employees":10000},
        {"id":6,"account_name":"Rogers Communications","industry":"Telecom","status":"Unknown","total_it_spend":28700000,"current_monthly_fee":0,"cdn_primary":"Akamai","security_primary":"Palo Alto Networks","billing_country":"CA","employees":25000},
    ],"total":2307,"page":1,"limit":50},
    '/api/lead-scores': [
        {"id":5,"account_name":"Shopify Inc","industry":"Computer Software","score":92,"total_it_spend":45200000,"current_monthly_fee":5400,"factors":[{"factor":"IT Spend","points":20,"detail":"$45.2M/mo"},{"factor":"Low Wallet Penetration","points":15,"detail":"0.01%"},{"factor":"Multi-Vendor Displacement","points":15,"detail":"3 vendors"}]},
        {"id":1,"account_name":"Apple Canada Inc","industry":"Telecom","score":85,"total_it_spend":35300000,"current_monthly_fee":0,"factors":[{"factor":"IT Spend","points":20,"detail":"$35.3M/mo"},{"factor":"Low Wallet Penetration","points":15,"detail":"0.0%"},{"factor":"Competitor Present","points":8,"detail":"Akamai"}]},
        {"id":6,"account_name":"Rogers Communications","industry":"Telecom","score":78,"total_it_spend":28700000,"current_monthly_fee":0,"factors":[{"factor":"IT Spend","points":20,"detail":"$28.7M/mo"},{"factor":"Low Wallet Penetration","points":15,"detail":"0.0%"},{"factor":"Enterprise","points":10,"detail":"25K emp"}]},
        {"id":2,"account_name":"Equinix Canada Ltd","industry":"Data Centers","score":72,"total_it_spend":22000000,"current_monthly_fee":0,"factors":[{"factor":"IT Spend","points":20,"detail":"$22.0M/mo"},{"factor":"Low Wallet Penetration","points":15,"detail":"0.0%"}]},
        {"id":3,"account_name":"Hotbot Limited","industry":"Publishing","score":65,"total_it_spend":16600000,"current_monthly_fee":0,"factors":[{"factor":"IT Spend","points":15,"detail":"$16.6M/mo"},{"factor":"Recent Activity","points":10,"detail":"85 days ago"}]},
    ],
    '/api/opportunities': [
        {"id":1,"account_name":"Shopify Inc","industry":"Computer Software","country":"CA","acv":1320000,"stage":"qualification","notes":"High IT spend, multi-vendor displacement","created_at":"2026-04-10"},
        {"id":2,"account_name":"Apple Canada Inc","industry":"Telecom","country":"CA","acv":576000,"stage":"prospecting","notes":"Massive Akamai displacement target","created_at":"2026-04-10"},
        {"id":3,"account_name":"Rogers Communications","industry":"Telecom","country":"CA","acv":420000,"stage":"prospecting","notes":"Enterprise consolidation play","created_at":"2026-04-10"},
    ],
    '/api/acv': {"totalAcv":2316000,"byCountry":[{"country":"CA","total":2316000}],"byStage":[{"stage":"qualification","cnt":1,"total":1320000},{"stage":"prospecting","cnt":2,"total":996000}]},
    '/api/alerts': [
        {"id":1,"alert_type":"infrastructure_change","title":"CDN Change: Shopify Inc","detail":"Fastly -> Cloudflare","severity":"high","read":0,"created_at":"2026-04-10T12:00:00Z"},
        {"id":2,"alert_type":"threat_match","title":"DDoS Attack: Telecom Sector","detail":"Major DDoS targeting Canadian telecoms","severity":"critical","read":0,"created_at":"2026-04-09T18:00:00Z"},
        {"id":3,"alert_type":"infrastructure_change","title":"DNS Change: Rogers","detail":"Route53 -> Cloudflare DNS","severity":"medium","read":1,"created_at":"2026-04-08T10:00:00Z"},
    ],
    '/api/team-stats': {"users":[{"user_email":"apotta@cloudflare.com"}],"emailsByUser":[{"user_email":"apotta@cloudflare.com","cnt":9}],"researchByUser":[{"user_email":"apotta@cloudflare.com","cnt":11}],"campaignsByUser":[{"user_email":"apotta@cloudflare.com","cnt":1}],"opportunitiesByUser":[{"user_email":"apotta@cloudflare.com","cnt":3,"total_acv":2316000}]},
    '/api/playbooks': [
        {"id":1,"name":"Competitive Displacement - Akamai","persona":"ae","industry":"","template":"## Akamai Displacement\n\n1. Acknowledge Akamai investment\n2. Highlight integrated platform\n3. Present CDN+Security POC","usage_count":5,"created_by":"apotta@cloudflare.com"},
        {"id":2,"name":"Zero Trust Initial Outreach","persona":"bdr","industry":"Financial Services","template":"## ZT Discovery\n\n1. Remote workforce size\n2. VPN pain points\n3. Compliance needs","usage_count":3,"created_by":"apotta@cloudflare.com"},
    ],
}

class Handler(http.server.SimpleHTTPRequestHandler):
    def __init__(self, *a, **kw):
        super().__init__(*a, directory=PUBLIC, **kw)
    def do_GET(self):
        path = self.path.split('?')[0]
        if path.startswith('/api/'):
            self.send_response(200)
            self.send_header('Content-Type','application/json')
            self.send_header('Access-Control-Allow-Origin','*')
            self.end_headers()
            self.wfile.write(json.dumps(MOCK.get(path, [])).encode())
            return
        if '.' not in path.split('/')[-1] and path != '/':
            self.path = '/index.html'
        super().do_GET()
    def do_POST(self):
        self.send_response(200)
        self.send_header('Content-Type','application/json')
        self.send_header('Access-Control-Allow-Origin','*')
        self.end_headers()
        self.wfile.write(b'{"success":true}')
    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin','*')
        self.send_header('Access-Control-Allow-Methods','GET,POST,DELETE,OPTIONS')
        self.send_header('Access-Control-Allow-Headers','Content-Type')
        self.end_headers()
    def log_message(self, *a): pass

PORT = 8789
server = http.server.HTTPServer(('127.0.0.1', PORT), Handler)
t = threading.Thread(target=server.serve_forever, daemon=True)
t.start()
print(f'[1/3] Mock server running on http://localhost:{PORT}')
time.sleep(1)

# ══════════════════════════════════════════════════════════════════
# STEP 2: Screenshots
# ══════════════════════════════════════════════════════════════════
print('[2/3] Capturing screenshots...')

from playwright.sync_api import sync_playwright

PAGES = [
    ('dashboard', f'http://localhost:{PORT}/#/', 2000),
    ('pipeline', f'http://localhost:{PORT}/#/pipeline', 2000),
    ('lead_scores', f'http://localhost:{PORT}/#/lead-scores', 2000),
    ('alerts', f'http://localhost:{PORT}/#/alerts', 1500),
    ('team', f'http://localhost:{PORT}/#/team', 1500),
    ('playbooks', f'http://localhost:{PORT}/#/playbooks', 1500),
]

with sync_playwright() as p:
    browser = p.chromium.launch(
        headless=True,
        executable_path='/opt/homebrew/lib/python3.14/site-packages/playwright/driver/package/.local-browsers/chromium_headless_shell-1208/chrome-headless-shell-mac-arm64/chrome-headless-shell'
    )
    ctx = browser.new_context(viewport={'width':1440,'height':900}, device_scale_factor=2, color_scheme='dark')
    for name, url, wait in PAGES:
        page = ctx.new_page()
        page.goto(url, wait_until='networkidle', timeout=15000)
        page.wait_for_timeout(wait)
        path = os.path.join(SHOTS, f'{name}.png')
        page.screenshot(path=path, full_page=False)
        print(f'  {name}.png')
        page.close()
    browser.close()

server.shutdown()
print(f'  Screenshots saved to {SHOTS}/')

# ══════════════════════════════════════════════════════════════════
# STEP 3: Build PPTX
# ══════════════════════════════════════════════════════════════════
print('[3/3] Building presentation...')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x08,0x09,0x0A)
SURF = RGBColor(0x16,0x17,0x18)
ELEV = RGBColor(0x1E,0x1F,0x21)
ACC = RGBColor(0x7C,0x7F,0xFF)
GRN = RGBColor(0x34,0xD3,0x99)
AMB = RGBColor(0xFB,0xBF,0x24)
RED = RGBColor(0xF8,0x71,0x71)
BLU = RGBColor(0x60,0xA5,0xFA)
PUR = RGBColor(0xA7,0x8B,0xFA)
ORG = RGBColor(0xFB,0x92,0x3C)
WHT = RGBColor(0xF7,0xF8,0xF8)
GRY = RGBColor(0xC4,0xC9,0xD4)
MUT = RGBColor(0x8A,0x8F,0x98)
DIM = RGBColor(0x5A,0x5E,0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def rect(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background(); return sh
def bar(s,l,t,w,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,Pt(3)); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def txt(s,l,t,w,h,text,sz=18,c=WHT,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='Calibri'; p.alignment=a; return tb
def stat(s,l,t,w,h,v,lb,c):
    rect(s,l,t,w,h,ELEV); bar(s,l,t,w,c); txt(s,l+Inches(0.2),t+Inches(0.2),w-Inches(0.4),Inches(0.5),v,sz=26,c=c,b=True); txt(s,l+Inches(0.2),t+Inches(0.7),w-Inches(0.4),Inches(0.3),lb,sz=10,c=MUT)
def img(s,path,l,t,max_w,max_h=None):
    """Add image preserving aspect ratio within max_w x max_h bounds."""
    if not os.path.exists(path): return
    from PIL import Image as PILImage
    with PILImage.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if max_h:
        # Fit within both max_w and max_h
        w_from_h = int(max_h * aspect)
        if w_from_h <= max_w:
            s.shapes.add_picture(path, l, t, w_from_h, max_h)
        else:
            h_from_w = int(max_w / aspect)
            s.shapes.add_picture(path, l, t, max_w, h_from_w)
    else:
        s.shapes.add_picture(path, l, t, width=max_w)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(1.5),Inches(2.2),Inches(2),ACC)
txt(s,Inches(1.5),Inches(2.5),Inches(10),Inches(1.2),'RevFlare',sz=72,c=WHT,b=True)
txt(s,Inches(1.5),Inches(3.8),Inches(8),Inches(0.8),'AI-Powered Sales Intelligence for Cloudflare',sz=30,c=ACC)
txt(s,Inches(1.5),Inches(5.0),Inches(8),Inches(1.0),'A single Cloudflare Worker powering deep account research, persona-curated outreach, competitive intelligence, threat monitoring, autonomous pipeline generation, MCP integrations, and compliant email execution.',sz=15,c=MUT)
txt(s,Inches(1.5),Inches(6.5),Inches(10),Inches(0.3),'Built by Arun Potta  |  Powered by Cloudflare Workers + AI',sz=12,c=DIM)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2: Business Value + ROI (merged)
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(1),Inches(0.6),Inches(1.5),ACC)
txt(s,Inches(1),Inches(0.75),Inches(10),Inches(0.6),'Business Value & ROI',sz=36,c=WHT,b=True)
txt(s,Inches(1),Inches(1.3),Inches(9),Inches(0.3),'RevFlare transforms how sales teams prospect, research, approve, and close deals.',sz=14,c=MUT)

# Value props - compact
vals = [
    ('10x Faster Research','8 live probes in parallel. 45 min -> 10 sec per account.',GRN),
    ('3x Higher Reply Rates','25 persona variants + AI chat refinement. Not templates.',BLU),
    ('100% Pipeline Coverage','AI agent auto-generates opps with realistic ACV.',PUR),
    ('Open Ecosystem','MCP client + server. Bring your own tools & agents.',AMB),
]
for i,(title,desc,c) in enumerate(vals):
    y=Inches(1.85+i*0.7)
    rect(s,Inches(1),y,Inches(5.5),Inches(0.55),SURF)
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1),y,Pt(3),Inches(0.55)); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    txt(s,Inches(1.3),y+Inches(0.05),Inches(2),Inches(0.22),title,sz=13,c=c,b=True)
    txt(s,Inches(1.3),y+Inches(0.28),Inches(5),Inches(0.22),desc,sz=10,c=GRY)

# Before/After - right side
txt(s,Inches(7),Inches(1.7),Inches(2.8),Inches(0.25),'WITHOUT',sz=10,c=RED,b=True)
txt(s,Inches(10),Inches(1.7),Inches(2.8),Inches(0.25),'WITH REVFLARE',sz=10,c=GRN,b=True)
befores = [('45 min','Research',RED),('2%','Reply rate',RED),('Manual','Pipeline',RED),('10+ tabs','Per prospect',RED)]
afters = [('10 sec','Research',GRN),('6-8%','Reply rate',GRN),('1 click','Pipeline',GRN),('1 platform','Unified',GRN)]
for i,(v,lb,c) in enumerate(befores):
    y=Inches(2.05+i*0.65)
    rect(s,Inches(7),y,Inches(2.7),Inches(0.5),SURF)
    txt(s,Inches(7.15),y+Inches(0.05),Inches(1),Inches(0.2),v,sz=13,c=c,b=True)
    txt(s,Inches(7.15),y+Inches(0.27),Inches(2.4),Inches(0.18),lb,sz=9,c=MUT)
for i,(v,lb,c) in enumerate(afters):
    y=Inches(2.05+i*0.65)
    rect(s,Inches(10),y,Inches(2.7),Inches(0.5),SURF)
    txt(s,Inches(10.15),y+Inches(0.05),Inches(1),Inches(0.2),v,sz=13,c=c,b=True)
    txt(s,Inches(10.15),y+Inches(0.27),Inches(2.4),Inches(0.18),lb,sz=9,c=MUT)

# Bottom ROI callout
rect(s,Inches(1),Inches(4.85),Inches(11.3),Inches(0.45),ELEV)
txt(s,Inches(1),Inches(4.88),Inches(11.3),Inches(0.4),'Annual value: 10 reps x 7.5 hrs/week saved x 50 weeks x $75/hr = $281,250 in recovered selling time',sz=12,c=GRN,a=PP_ALIGN.CENTER)

# Architecture mini at bottom
txt(s,Inches(1),Inches(5.6),Inches(5),Inches(0.25),'CLOUDFLARE EDGE ARCHITECTURE',sz=10,c=DIM,b=True)
svcs = [('Workers','App server',ACC),('D1','29 tables',GRN),('Workers AI','3 LLMs',PUR),('Browser','Scraping',BLU),('KV','Cache',AMB),('Cron','Auto-send',RED)]
for i,(nm,ds,c) in enumerate(svcs):
    x=Inches(1+i*2.0); y=Inches(5.95)
    rect(s,x,y,Inches(1.8),Inches(1.15),SURF); bar(s,x,y,Inches(1.8),c)
    txt(s,x+Inches(0.12),y+Inches(0.2),Inches(1.6),Inches(0.25),nm,sz=12,c=c,b=True)
    txt(s,x+Inches(0.12),y+Inches(0.5),Inches(1.6),Inches(0.5),ds,sz=9,c=MUT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3: Dashboard Screenshot
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(0.6),Inches(0.4),Inches(1.5),ACC)
txt(s,Inches(0.6),Inches(0.55),Inches(8),Inches(0.5),'Account Intelligence Dashboard',sz=28,c=WHT,b=True)
txt(s,Inches(0.6),Inches(1.1),Inches(8),Inches(0.3),'2,307 accounts with IT spend, competitor stack, and engagement signals at a glance.',sz=13,c=MUT)
shot = os.path.join(SHOTS,'dashboard.png')
if os.path.exists(shot):
    rect(s,Inches(0.55),Inches(1.5),Inches(12.2),Inches(5.75),RGBColor(0x0C,0x0D,0x0E))
    img(s,shot,Inches(0.6),Inches(1.55),Inches(12.1),Inches(5.65))
else:
    txt(s,Inches(3),Inches(3.5),Inches(7),Inches(0.5),'[Dashboard Screenshot]',sz=20,c=DIM,a=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4: Pipeline + Lead Scores (side-by-side screenshots)
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(0.6),Inches(0.4),Inches(1.5),ACC)
txt(s,Inches(0.6),Inches(0.55),Inches(5),Inches(0.4),'Pipeline & AI Lead Scoring',sz=28,c=WHT,b=True)
txt(s,Inches(0.6),Inches(1.0),Inches(12),Inches(0.3),'Left: AI Opportunity Agent auto-generates deals with ACV.  Right: 8-factor lead scoring ranks every account 0-100.',sz=12,c=MUT)
shot_p = os.path.join(SHOTS,'pipeline.png')
shot_l = os.path.join(SHOTS,'lead_scores.png')
if os.path.exists(shot_p):
    rect(s,Inches(0.35),Inches(1.4),Inches(6.35),Inches(5.8),RGBColor(0x0C,0x0D,0x0E))
    img(s,shot_p,Inches(0.4),Inches(1.45),Inches(6.25),Inches(5.7))
if os.path.exists(shot_l):
    rect(s,Inches(6.85),Inches(1.4),Inches(6.25),Inches(5.8),RGBColor(0x0C,0x0D,0x0E))
    img(s,shot_l,Inches(6.9),Inches(1.45),Inches(6.15),Inches(5.7))

# ══════════════════════════════════════════════════════════════════
# SLIDE 5: 20 Features + Competitive Intel (merged)
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(0.5),Inches(0.5),Inches(1.5),ACC)
txt(s,Inches(0.5),Inches(0.65),Inches(10),Inches(0.5),'20 Features + Competitive Intelligence',sz=32,c=WHT,b=True)

# Features - 5x4 compact grid (20 features)
feats = [
    ('Lead Scoring','AI 0-100',GRN),('ROI Calc','Savings est.',BLU),('Lookalikes','Match stack',PUR),('Meeting Prep','Call brief',AMB),('Sequences','Auto-execute',ORG),
    ('Change Detect','Infra scan',RED),('A/B Testing','2 variants',ACC),('Voice Notes','To email',GRN),('Alerts','+ AI products',RED),('Team Dash','Leaderboard',BLU),
    ('Playbooks','Templates',PUR),('Search','AI-ranked',AMB),('Pipeline','ACV track',GRN),('Win/Loss','Analysis',ORG),('Opp Agent','Auto gen',ACC),
    ('Email Approval','Review & gate',GRN),('AI Chat Refine','Rewrite in place',BLU),('Open Tracking','Pixel + stats',PUR),('CAN-SPAM','Unsubscribe',AMB),('MCP Gateway','Client + server',ORG),
]
for i,(nm,ds,c) in enumerate(feats):
    col=i%5; row=i//5; x=Inches(0.5+col*1.65); y=Inches(1.35+row*0.75)
    rect(s,x,y,Inches(1.5),Inches(0.62),SURF); bar(s,x,y,Inches(1.5),c)
    txt(s,x+Inches(0.08),y+Inches(0.06),Inches(1.35),Inches(0.22),nm,sz=10,c=c,b=True)
    txt(s,x+Inches(0.08),y+Inches(0.28),Inches(1.35),Inches(0.3),ds,sz=8,c=MUT)

# Competitive - right half
txt(s,Inches(8.75),Inches(1.25),Inches(4),Inches(0.25),'COMPETITIVE (12 CATEGORIES, 40+ RIVALS)',sz=9,c=DIM,b=True)
cats = ['CDN -- Akamai, CloudFront, Fastly','WAF -- Imperva, AWS WAF, F5','DDoS -- Prolexic, AWS Shield','Zero Trust -- Zscaler, Palo Alto','Edge Compute -- Lambda@Edge, Vercel','DNS -- Route 53, NS1, Google']
for i,ct in enumerate(cats):
    y=Inches(1.6+i*0.42); rect(s,Inches(8.75),y,Inches(4.1),Inches(0.32),SURF)
    txt(s,Inches(8.88),y+Inches(0.04),Inches(3.85),Inches(0.22),ct,sz=9,c=GRY)

# Threat Intel - bottom
txt(s,Inches(0.5),Inches(4.55),Inches(5),Inches(0.25),'THREAT INTELLIGENCE',sz=9,c=DIM,b=True)
thr = [('26 RSS Feeds','Real-time incidents',RED),('GDELT + 5 News APIs','Google, Bing, NewsAPI, GNews, MediaStack',AMB),('Auto-Match','By industry & country',BLU),('Trigger Emails','Incident outreach',GRN),('Nightly Cron','Daily scan + alerts',PUR)]
for i,(t,d,c) in enumerate(thr):
    x=Inches(0.5+i*2.55); y=Inches(4.85)
    rect(s,x,y,Inches(2.4),Inches(0.6),SURF)
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,Pt(3),Inches(0.6)); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    txt(s,x+Inches(0.15),y+Inches(0.04),Inches(2.1),Inches(0.2),t,sz=10,c=c,b=True)
    txt(s,x+Inches(0.15),y+Inches(0.28),Inches(2.1),Inches(0.2),d,sz=8,c=MUT)

# Security + Integrations - bottom strip
txt(s,Inches(0.5),Inches(5.65),Inches(5),Inches(0.25),'SECURITY & INTEGRATIONS (47 FIXES ACROSS 3 AUDIT ROUNDS)',sz=9,c=DIM,b=True)
items = [('Access JWT',GRN),('AES-256-GCM',BLU),('SSRF Block',RED),('XSS Escaped',AMB),('SOQL Safe',PUR),('CORS Locked',ACC),('Gmail OAuth',GRN),('SF OAuth',BLU),('Share Links',PUR),('MCP Secure',ORG)]
for i,(t,c) in enumerate(items):
    x=Inches(0.5+i*1.27); y=Inches(5.95)
    rect(s,x,y,Inches(1.2),Inches(0.4),SURF); bar(s,x,y,Inches(1.2),c)
    txt(s,x+Inches(0.08),y+Inches(0.08),Inches(1.1),Inches(0.25),t,sz=8,c=c,b=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6: Alerts + Team + Playbooks screenshots
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(0.6),Inches(0.4),Inches(1.5),ACC)
txt(s,Inches(0.6),Inches(0.55),Inches(5),Inches(0.4),'Alerts, Team & Playbooks',sz=28,c=WHT,b=True)
txt(s,Inches(0.6),Inches(1.0),Inches(12),Inches(0.3),'Real-time infrastructure alerts, team activity leaderboard, and reusable sales playbooks.',sz=12,c=MUT)

shot_a = os.path.join(SHOTS,'alerts.png')
shot_t = os.path.join(SHOTS,'team.png')
shot_pb = os.path.join(SHOTS,'playbooks.png')
if os.path.exists(shot_a):
    rect(s,Inches(0.35),Inches(1.4),Inches(4.2),Inches(5.8),RGBColor(0x0C,0x0D,0x0E))
    img(s,shot_a,Inches(0.4),Inches(1.45),Inches(4.1),Inches(5.7))
if os.path.exists(shot_t):
    rect(s,Inches(4.7),Inches(1.4),Inches(4.2),Inches(5.8),RGBColor(0x0C,0x0D,0x0E))
    img(s,shot_t,Inches(4.75),Inches(1.45),Inches(4.1),Inches(5.7))
if os.path.exists(shot_pb):
    rect(s,Inches(9.05),Inches(1.4),Inches(4.05),Inches(5.8),RGBColor(0x0C,0x0D,0x0E))
    img(s,shot_pb,Inches(9.1),Inches(1.45),Inches(3.95),Inches(5.7))

# ══════════════════════════════════════════════════════════════════
# SLIDE 7: MCP + Email Tracking + AI Chat (what's new)
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s,Inches(0.5),Inches(0.4),Inches(1.5),ACC)
txt(s,Inches(0.5),Inches(0.55),Inches(12),Inches(0.5),'Open Ecosystem, Compliant Outreach, Conversational AI',sz=28,c=WHT,b=True)
txt(s,Inches(0.5),Inches(1.1),Inches(12),Inches(0.3),'Four new capabilities bringing RevFlare from standalone tool to enterprise platform.',sz=12,c=MUT)

# ── MCP Integration (top-left, full height) ────────────────────────
rect(s,Inches(0.5),Inches(1.6),Inches(6.2),Inches(2.75),ELEV); bar(s,Inches(0.5),Inches(1.6),Inches(6.2),ORG)
txt(s,Inches(0.7),Inches(1.7),Inches(5.8),Inches(0.3),'MCP (Model Context Protocol)',sz=14,c=ORG,b=True)
txt(s,Inches(0.7),Inches(2.0),Inches(5.8),Inches(0.3),'Bi-directional. Client + Server.',sz=10,c=MUT)

txt(s,Inches(0.7),Inches(2.35),Inches(2.8),Inches(0.22),'AS CLIENT',sz=9,c=DIM,b=True)
mcp_client = ['+ Netstrat (network + strategy)','+ Google Workspace (calendar, contacts)','+ Wiki (internal knowledge)','+ Jira (issue context)','+ Cloudflare Docs (product lookup)','+ SSRF-hardened URL validator']
for i,item in enumerate(mcp_client):
    txt(s,Inches(0.7),Inches(2.55+i*0.22),Inches(3.0),Inches(0.2),item,sz=8,c=GRY)

txt(s,Inches(3.8),Inches(2.35),Inches(2.8),Inches(0.22),'AS SERVER',sz=9,c=DIM,b=True)
mcp_server = ['+ lookup_account','+ get_lead_score','+ get_account_research','+ get_pipeline','+ get_alerts','+ get_email_stats']
for i,item in enumerate(mcp_server):
    txt(s,Inches(3.8),Inches(2.55+i*0.22),Inches(3.0),Inches(0.2),item,sz=8,c=GRY)

# ── AI Chat Email Refinement (top-right) ───────────────────────────
rect(s,Inches(6.85),Inches(1.6),Inches(6.2),Inches(2.75),ELEV); bar(s,Inches(6.85),Inches(1.6),Inches(6.2),BLU)
txt(s,Inches(7.05),Inches(1.7),Inches(5.8),Inches(0.3),'AI Chat Email Refinement',sz=14,c=BLU,b=True)
txt(s,Inches(7.05),Inches(2.0),Inches(5.8),Inches(0.3),'Conversational rewrites. Full turn history.',sz=10,c=MUT)

chat_feats = [
    '+ Plain-English instructions: "make it shorter, lead with ROI"',
    '+ Llama 3.3 70B multi-turn with context memory (20 prior turns)',
    '+ Rewrites in place, resets to pending_approval for re-review',
    '+ Full chat history persisted per message for auditability',
    '+ Works on persona messages AND campaign emails',
    '+ Subject line auto-extracted and updated alongside body',
]
for i,item in enumerate(chat_feats):
    txt(s,Inches(7.05),Inches(2.4+i*0.28),Inches(6.0),Inches(0.25),item,sz=9,c=GRY)

# ── Email Tracking & CAN-SPAM (bottom-left) ────────────────────────
rect(s,Inches(0.5),Inches(4.5),Inches(6.2),Inches(2.75),ELEV); bar(s,Inches(0.5),Inches(4.5),Inches(6.2),PUR)
txt(s,Inches(0.7),Inches(4.6),Inches(5.8),Inches(0.3),'Email Tracking & CAN-SPAM',sz=14,c=PUR,b=True)
txt(s,Inches(0.7),Inches(4.9),Inches(5.8),Inches(0.3),'Opens, unsubscribes, suppression - built in.',sz=10,c=MUT)

track_feats = [
    '+ 1x1 tracking pixel per email (unique tracking_id)',
    '+ One-click List-Unsubscribe header + landing page',
    '+ Per-user suppression list (bounces, complaints)',
    '+ Auto-block of suppressed addresses before send',
    '+ Public endpoints bypass Access for deliverability',
    '+ Gmail signature auto-appended to every outbound',
]
for i,item in enumerate(track_feats):
    txt(s,Inches(0.7),Inches(5.3+i*0.28),Inches(6.0),Inches(0.25),item,sz=9,c=GRY)

# ── Email Performance Dashboard (bottom-right) ─────────────────────
rect(s,Inches(6.85),Inches(4.5),Inches(6.2),Inches(2.75),ELEV); bar(s,Inches(6.85),Inches(4.5),Inches(6.2),GRN)
txt(s,Inches(7.05),Inches(4.6),Inches(5.8),Inches(0.3),'Email Performance Dashboard',sz=14,c=GRN,b=True)
txt(s,Inches(7.05),Inches(4.9),Inches(5.8),Inches(0.3),'Funnel, trends, suppression - one screen.',sz=10,c=MUT)

perf_feats = [
    '+ Sent / opened / replied rates per user',
    '+ Daily send trend (last 14 days) with cap tracking',
    '+ Per-campaign funnel (last 20 campaigns)',
    '+ Suppression list management with remove button',
    '+ Today\'s usage vs 100/day cap (reputation protection)',
    '+ Per-message open count + reply flag on every email',
]
for i,item in enumerate(perf_feats):
    txt(s,Inches(7.05),Inches(5.3+i*0.28),Inches(6.0),Inches(0.25),item,sz=9,c=GRY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8: Closing
# ══════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
gl=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3),Inches(0),Inches(7),Inches(5)); gl.fill.solid(); gl.fill.fore_color.rgb=RGBColor(0x1A,0x1D,0x3A); gl.line.fill.background()

bar(s,Inches(1),Inches(0.8),Inches(1.5),ACC)
txt(s,Inches(1),Inches(1.0),Inches(10),Inches(0.6),'RevFlare by the Numbers',sz=40,c=WHT,b=True)

nums = [('~12,700','Lines of Code',GRN),('29','DB Tables',BLU),('108','API Endpoints',PUR),('8','Live Probes',AMB),('25','Email Variants',ORG),
        ('12','Competitive Cats',RED),('40+','Competitors',ACC),('26','Threat Feeds',RED),('5','Personas',GRN),('$0','Infra Cost',AMB)]
for i,(v,lb,c) in enumerate(nums):
    col=i%5; row=i//5; x=Inches(0.7+col*2.45); y=Inches(2.0+row*1.5)
    stat(s,x,y,Inches(2.2),Inches(1.1),v,lb,c)

rect(s,Inches(2.5),Inches(5.4),Inches(8.3),Inches(1.6),SURF)
txt(s,Inches(2.5),Inches(5.5),Inches(8.3),Inches(0.5),'revflare.arunpotta1024.workers.dev',sz=22,c=ACC,b=True,a=PP_ALIGN.CENTER)
txt(s,Inches(2.5),Inches(6.0),Inches(8.3),Inches(0.35),'github.com/pottaarun/RevFlare',sz=14,c=MUT,a=PP_ALIGN.CENTER)
txt(s,Inches(2.5),Inches(6.4),Inches(8.3),Inches(0.3),'Built by Arun Potta  |  Powered by Cloudflare Workers + AI',sz=12,c=DIM,a=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────
out = os.path.join(ROOT, 'RevFlare-Presentation.pptx')
prs.save(out)
print(f'\nSaved: {out}')
print(f'Slides: {len(prs.slides)}')
shots_found = sum(1 for f in os.listdir(SHOTS) if f.endswith('.png'))
print(f'Screenshots embedded: {shots_found}')
